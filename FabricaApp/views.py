from .models import (
    FormularioProyectoInterno,
    FormularioProyectoFabrica,
    FormularioProyectoFabLab,
    FabLabImage,
    FabricaImage,
)
from .forms import (
    ProyectoInternoCreateForm,
    ProyectoFabricaCreateForm,
    ProyectoFabLabCreateForm,
    ProyectoFabricaFondosCreateForm,
    ImageForm,
    ImageFabricaForm,
)
from django.contrib.auth.mixins import LoginRequiredMixin
from django.views.generic import (
    CreateView,
    ListView,
    DetailView,
    DeleteView,
    UpdateView,
    View,
)
from django.contrib import messages
from django.urls import reverse_lazy
from django.shortcuts import redirect, get_object_or_404, render
from django.db.models import Q
from django.core.paginator import Paginator
from core.mixins import PermitsPositionMixin
from django.http import HttpResponse

# Create your views here.


##########################################
##### PROYECTO FABRICA INTERNO CARLA #####
##########################################


class ProyectoInternoCreateView(LoginRequiredMixin, PermitsPositionMixin, CreateView):
    model = FormularioProyectoInterno
    form_class = ProyectoInternoCreateForm
    template_name = "pages/fabrica/fabrica.html"
    success_url = reverse_lazy("FabriCreate")

    def form_valid(self, form):
        user = self.request.user
        formulario_proyecto = form.save(commit=False)
        formulario_proyecto.user_id = user
        formulario_proyecto.save()
        messages.success(self.request, "Archivo cargado con éxito")
        return super().form_valid(form)

    def form_invalid(self, form):
        for field, errors in form.errors.items():
            for error in errors:
                print(field, error)
                messages.error(self.request, f"{error}")
        return redirect("FabriCreate")


class ProyectoInternoListView(LoginRequiredMixin, PermitsPositionMixin, ListView):
    model = FormularioProyectoInterno
    template_name = "pages/fabrica/fabrica_lista.html"
    paginate_by = 8

    def get_queryset(self):
        queryset = super().get_queryset().order_by("-id")
        search_query = self.request.GET.get("search")

        if search_query:
            queryset = queryset.filter(
                Q(docente_id__run=search_query)
                | Q(docente_id__nombre__icontains=search_query)
                | Q(docente_id__apellido_p__icontains=search_query)
                | Q(docente_id__email__icontains=search_query)
            )
        return queryset

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        paginator = Paginator(context["object_list"], self.paginate_by)
        page = self.request.GET.get("page")
        context["object_list"] = paginator.get_page(page)
        context["placeholder"] = "Buscar por run, nombre, apellido o correo."
        return context


class ProyectoInternoDetailView(LoginRequiredMixin, PermitsPositionMixin, DetailView):
    model = FormularioProyectoInterno
    template_name = "pages/fabrica/fabrica_detalle.html"
    context_object_name = "item"

    def get_object(self):
        id_ = self.kwargs.get("id")
        return get_object_or_404(FormularioProyectoInterno, id=id_)


class ProyectoInternoDeleteView(LoginRequiredMixin, PermitsPositionMixin, DeleteView):
    model = FormularioProyectoInterno
    success_url = reverse_lazy("FabriList")

    def get(self, request, *args, **kwargs):
        self.object = self.get_object()
        messages.success(self.request, "Proyecto eliminado correctamente")
        self.object.delete()
        return redirect(self.get_success_url())


class ProyectoInternoUpdateView(LoginRequiredMixin, PermitsPositionMixin, UpdateView):
    model = FormularioProyectoInterno
    form_class = ProyectoInternoCreateForm
    template_name = "pages/fabrica/fabrica.html"
    success_url = reverse_lazy("FabriList")

    def form_valid(self, form):
        form.clean()
        form.save()
        messages.success(self.request, "Proyecto editado correctamente")
        return super().form_valid(form)

    def form_invalid(self, form):
        messages.error(self.request, "Error en el formulario")
        for field, errors in form.errors.items():
            for error in errors:
                messages.error(self.request, f"{error}")
        return redirect("FabriList")


######################################
### ### PROYECTO FABRICA ANITA ### ###
######################################


class ProyectoFabricaCreateView(LoginRequiredMixin, View):

    def handle_form_errors(self, proyecto_form, img_form, fondos_form):
        combined_errors = {}

        errors = [proyecto_form, img_form]
        if proyecto_form.cleaned_data.get("fondos", False):
            errors.append(fondos_form)

        # Combinar errores de todos los formularios
        for form in errors:
            for field, field_errors in form.errors.items():
                combined_errors.setdefault(field, []).extend(field_errors)

        # Registrar los mensajes de error
        for field, field_errors in combined_errors.items():
            for error in field_errors:
                messages.error(self.request, error)

    def get(self, request, *args, **kwargs):
        parte1_form = ProyectoFabricaCreateForm()
        parte2_form = ProyectoFabricaFondosCreateForm()
        image_form = ImageFabricaForm()
        return render(
            request,
            "pages/ficha/fabrica_fondos.html",
            {
                "parte1_form": parte1_form,
                "image_form": image_form,
                "parte2_form": parte2_form,
            },
        )

    def post(self, request, *args, **kwargs):
        parte1_form = ProyectoFabricaCreateForm(request.POST, request.FILES)
        image_form = ImageFabricaForm(request.POST, request.FILES)

        if parte1_form.is_valid() and image_form.is_valid():
            user = self.request.user
            instance = parte1_form.save(commit=False)
            instance.user_id = user
            instance.save()

            uploaded_images = request.FILES.getlist("image")
            for image in uploaded_images:
                FabricaImage.objects.create(ficha_fabrica=instance, image=image)

            messages.success(self.request, "Creado con éxito")

            if parte1_form.cleaned_data["fondos"]:
                parte2_form = ProyectoFabricaFondosCreateForm(request.POST)
                if parte2_form.is_valid():
                    fondos_instance = parte2_form.save(commit=False)
                    fondos_instance.proyecto = instance
                    fondos_instance.save()
                    return redirect("FabriFichaCreate")
            else:
                return redirect("FabriFichaCreate")
        else:
            parte2_form = (
                ProyectoFabricaFondosCreateForm()
            )  # Crea un formulario vacío si parte1_form no es válido
            self.handle_form_errors(parte1_form, image_form, parte2_form)

        return render(
            request,
            "pages/ficha/fabrica_fondos.html",
            {
                "parte1_form": parte1_form,
                "image_form": image_form,
                "parte2_form": parte2_form,
            },
        )


class ProyectoFabricaListView(LoginRequiredMixin, ListView):
    model = FormularioProyectoFabrica
    template_name = "pages/ficha/fabrica_lista.html"
    paginate_by = 8

    def get_queryset(self):
        queryset = super().get_queryset().order_by("-id")
        user = self.request.user
        search_query = self.request.GET.get("search")

        # Filtrar según el usuario
        if not user.is_superuser:
            queryset = queryset.filter(user_id=user.id)

        # Filtrar por código SIR si se ha proporcionado una búsqueda
        if search_query:
            queryset = queryset.filter(Q(codigo_sir=search_query))

        return queryset

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        paginator = Paginator(context["object_list"], self.paginate_by)
        page = self.request.GET.get("page")
        context["object_list"] = paginator.get_page(page)
        context["placeholder"] = "Buscar por código SIR"
        return context


class ProyectoFabricaDetailView(LoginRequiredMixin, DetailView):
    model = FormularioProyectoFabrica
    template_name = "pages/ficha/fabrica_detalle.html"
    context_object_name = "item"

    def get_object(self):
        id_ = self.kwargs.get("id")
        return get_object_or_404(self.model, id=id_)


class ProyectoFabricaUpdateView(LoginRequiredMixin, View):
    template_name = "pages/ficha/fabrica_update.html"

    def handle_form_errors(self, proyecto_form, fondos_form):
        combined_errors = {}

        errors = [proyecto_form]
        if proyecto_form.cleaned_data.get("fondos", False):
            errors.append(fondos_form)

        # Combinar errores de todos los formularios
        for form in errors:
            for field, field_errors in form.errors.items():
                combined_errors.setdefault(field, []).extend(field_errors)

        # Registrar los mensajes de error
        for field, field_errors in combined_errors.items():
            for error in field_errors:
                messages.error(self.request, error)

    def get(self, request, *args, **kwargs):
        id_ = self.kwargs.get("pk")
        proyecto = get_object_or_404(FormularioProyectoFabrica, id=id_)

        proyecto_form = ProyectoFabricaCreateForm(instance=proyecto)

        fondos_form = ProyectoFabricaFondosCreateForm(
            instance=(
                proyecto.fondos_proyecto.first()
                if proyecto.fondos_proyecto.exists()
                else None
            )
        )

        fecha_inicio = (
            proyecto.fecha_inicio.strftime("%Y-%m-%d") if proyecto.fecha_inicio else ""
        )

        context = {
            "proyecto_form": proyecto_form,
            "fondos_form": fondos_form,
            "fecha_inicio": fecha_inicio,
        }

        return render(request, self.template_name, context)

    def post(self, request, *args, **kwargs):
        id_ = self.kwargs.get("pk")
        proyecto = get_object_or_404(FormularioProyectoFabrica, id=id_)

        proyecto_form = ProyectoFabricaCreateForm(
            request.POST, request.FILES, instance=proyecto
        )

        fondos_form = ProyectoFabricaFondosCreateForm(
            request.POST,
            instance=proyecto.fondos_proyecto.first() if proyecto.fondos else None,
        )

        if proyecto_form.is_valid():
            proyecto_form.save()

            if proyecto.fondos and fondos_form.is_valid():
                fondos = fondos_form.save(commit=False)
                fondos.proyecto = proyecto
                fondos.save()

            messages.success(request, "Proyecto actualizado exitosamente.")
            return redirect("FabriFichaList")
        else:
            self.handle_form_errors(proyecto_form, fondos_form)

        context = {"proyecto_form": proyecto_form, "fondos_form": fondos_form}
        return render(request, self.template_name, context)


class ProyectoFabricaDeleteView(LoginRequiredMixin, PermitsPositionMixin, DeleteView):
    model = FormularioProyectoFabrica
    success_url = reverse_lazy("FabriFichaList")

    def get(self, request, *args, **kwargs):
        self.object = self.get_object()
        messages.success(self.request, "Proyecto eliminado correctamente")
        self.object.delete()
        return redirect(self.get_success_url())


######################################
### ### PROYECTO FABLAB MIGUEL ### ###
######################################


class ProyectoFabLabCreateView(LoginRequiredMixin, PermitsPositionMixin, View):
    def get(self, request):
        fablab_form = ProyectoFabLabCreateForm()
        image_form = ImageForm()
        return render(
            request,
            "pages/fablab/fablab.html",
            {"fablab_form": fablab_form, "image_form": image_form},
        )

    def post(self, request):
        fablab_form = ProyectoFabLabCreateForm(request.POST)
        image_form = ImageForm(request.POST, request.FILES)

        if fablab_form.is_valid() and image_form.is_valid():
            user = self.request.user
            fablab = fablab_form.save(commit=False)
            fablab.user_id = user
            fablab.save()
            messages.success(self.request, "Creado con éxito")

            docentes = fablab_form.cleaned_data.get("docentes")
            fablab.docentes.set(docentes)

            uploaded_images = request.FILES.getlist("image")
            for image in uploaded_images:
                FabLabImage.objects.create(ficha_fablab=fablab, image=image)

            return redirect("FabLabFichaCreate")

        return render(
            request,
            "pages/fablab/fablab.html",
            {"fablab_form": fablab_form, "image_form": image_form},
        )


class ProyectoFabLabListView(LoginRequiredMixin, PermitsPositionMixin, ListView):
    model = FormularioProyectoFabLab
    template_name = "pages/fablab/fablab_lista.html"
    paginate_by = 8

    def get_queryset(self):
        queryset = super().get_queryset().order_by("-id")
        search_query = self.request.GET.get("search")

        if search_query:
            queryset = queryset.filter(Q(nombre_propuesta=search_query))
        return queryset

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        paginator = Paginator(context["object_list"], self.paginate_by)
        page = self.request.GET.get("page")
        context["object_list"] = paginator.get_page(page)
        context["placeholder"] = "Buscar por nombre de propuesta."
        return context


class ProyectoFabLabDetailView(LoginRequiredMixin, PermitsPositionMixin, DetailView):
    model = FormularioProyectoFabLab
    template_name = "pages/fablab/fablab_detalle.html"
    context_object_name = "item"

    def get_object(self):
        id_ = self.kwargs.get("id")
        return get_object_or_404(self.model, id=id_)


class ProyectoFabLabDeleteView(LoginRequiredMixin, PermitsPositionMixin, DeleteView):
    model = FormularioProyectoFabLab
    success_url = reverse_lazy("FabLabFichaList")

    def get(self, request, *args, **kwargs):
        self.object = self.get_object()
        messages.success(self.request, "Proyecto eliminado correctamente")
        self.object.delete()
        return redirect(self.get_success_url())


class ProyectoFabLAbUpdateView(LoginRequiredMixin, PermitsPositionMixin, UpdateView):
    model = FormularioProyectoFabLab
    form_class = ProyectoFabLabCreateForm
    template_name = "pages/fablab/fablab_update.html"
    success_url = reverse_lazy("FabLabFichaList")

    def form_valid(self, form):
        form.clean()
        form.save()
        messages.success(self.request, "Proyecto editado correctamente")
        return super().form_valid(form)

    def form_invalid(self, form):
        messages.error(self.request, "Error en el formulario")
        for field, errors in form.errors.items():
            for error in errors:
                messages.error(self.request, f"{error}")
        return redirect("FabLabFichaList")


############################
### ###  Generate PDF
############################
from io import BytesIO
from django.conf import settings
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Image,
    Table,
    TableStyle,
    PageBreak,
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.colors import red, HexColor
from reportlab.graphics.shapes import Drawing, Line
import os


class PdfView(View):
    def generate_pdf(self, formulario_proyecto):
        response = HttpResponse(content_type="application/pdf")
        response["Content-Disposition"] = 'attachment; filename="reporte.pdf"'
        doc = SimpleDocTemplate(
            response,
            pagesize=letter,
            leftMargin=20,
            rightMargin=20,
            topMargin=30,
            bottomMargin=20,
        )

        # Estilos
        styles = getSampleStyleSheet()
        normal_style = styles["Normal"]
        title_style = styles["Title"]
        bold_style = ParagraphStyle(
            name="Bold", parent=styles["Normal"], fontName="Helvetica-Bold"
        )
        red_bold_style = ParagraphStyle(
            name="RedBold",
            parent=styles["Normal"],
            fontName="Helvetica-Bold",
            fontSize=14,
            textColor=red,
        )
        dark_gray_color = HexColor("#333333")

        # Crear estilos personalizados para los párrafos

        normal_style.fontSize = 12
        normal_style.leading = 20
        normal_style.textColor = dark_gray_color
        normal_style.spaceAfter = 12

        # Página de portada
        story = []

        ############################################################################################
        # Imágenes de portada desde static
        static_path = settings.STATICFILES_DIRS[0]
        img_inacap_path = os.path.join(static_path, "img/inacap.jpg")
        img_logo_fabrica_path = os.path.join(static_path, "img/logo_fabrica.jpeg")

        img_inacap = (
            Image(img_inacap_path, 2 * inch, 1 * inch)
            if os.path.exists(img_inacap_path)
            else None
        )
        img_logo_fabrica = (
            Image(img_logo_fabrica_path, 1.6 * inch, 1.6 * inch)
            if os.path.exists(img_logo_fabrica_path)
            else None
        )

        table_data = []
        if img_inacap or img_logo_fabrica:
            row = []
            row.append(img_inacap)
            row.append(Spacer(4 * inch, 1 * inch))
            row.append(img_logo_fabrica)
            table_data.append(row)

        if table_data:
            table = Table(table_data, colWidths=[2 * inch, 4 * inch, 2 * inch])
            table.setStyle(
                TableStyle(
                    [
                        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ]
                )
            )
            story.append(table)

        #################################################################################################
        # Portada
        story.append(Spacer(1, 1 * inch))
        story.append(
            Paragraph(
                "Reporte de Proyecto: %s" % formulario_proyecto.nombre_propuesta,
                title_style,
            )
        )
        story.append(Spacer(1, 200))

        # Fecha y empresa (más abajo)
        story.append(
            Paragraph(
                "<strong>Fecha de Inicio:</strong> %s"
                % formulario_proyecto.fecha_inicio,
                bold_style,
            )
        )
        story.append(Spacer(1, 6))
        story.append(
            Paragraph(
                "<strong>Empresa:</strong> %s"
                % formulario_proyecto.empresa_id.nombre_empresa,
                bold_style,
            )
        )
        story.append(Spacer(1, 12))

        story.append(PageBreak())
        ###################################################################################################
        # Nueva página para el contenido

        story.append(
            Paragraph("<strong>Detalles del Proyecto</strong>", red_bold_style)
        )
        drawing = Drawing()
        line_width = 550
        line = Line(0, 0, line_width, 0)
        line.strokeColor = red
        line.strokeWidth = 1
        drawing.add(line)
        drawing.height = 10
        drawing.width = line_width
        story.append(drawing)
        ###################################################################################################

        story.append(Spacer(1, 12))

        # Añadir párrafos con separación y colores
        story.append(
            Paragraph(
                "<strong>Empresa:</strong> <font color='%s'>%s</font>"
                % (dark_gray_color, formulario_proyecto.empresa_id.nombre_empresa),
                normal_style,
            )
        )
        story.append(Spacer(1, 12))
        story.append(
            Paragraph(
                "<strong>Docente:</strong> <font color='%s'>%s %s</font>"
                % (
                    dark_gray_color,
                    formulario_proyecto.docente_id.nombre,
                    formulario_proyecto.docente_id.apellido_p,
                ),
                normal_style,
            )
        )
        story.append(Spacer(1, 12))
        story.append(
            Paragraph(
                "<strong>Problema/Desafío:</strong> <font color='%s'>%s</font>"
                % (dark_gray_color, formulario_proyecto.problema),
                normal_style,
            )
        )
        story.append(Spacer(1, 12))
        story.append(
            Paragraph(
                "<strong>Objetivo:</strong> <font color='%s'>%s</font>"
                % (dark_gray_color, formulario_proyecto.objetivo),
                normal_style,
            )
        )
        story.append(Spacer(1, 12))

        story.append(PageBreak())

        #####################################################################################################
        # Nueva página para el contenido

        story.append(
            Paragraph("<strong>Imágenes del Proyecto</strong>", red_bold_style)
        )
        drawing = Drawing()
        line_width = 550
        line = Line(0, 0, line_width, 0)
        line.strokeColor = red
        line.strokeWidth = 1
        drawing.add(line)
        drawing.height = 10
        drawing.width = line_width
        story.append(drawing)
        ###################################################################################################

        story.append(Spacer(1, 24))

        images = [
            Image(image.image.path, 3 * inch, 3 * inch)
            for image in formulario_proyecto.images.all()
        ]  # Tamaño aumentado

        # Crear tabla para las imágenes en dos columnas
        if images:
            image_table_data = []
            for i in range(0, len(images), 2):
                row = images[i : i + 2]
                image_table_data.append(row)

            image_table = Table(
                image_table_data, colWidths=[3.5 * inch] * 2
            )  # Ajustar el ancho de las columnas
            image_table.setStyle(
                TableStyle(
                    [
                        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        ("SPACEBEFORE", (0, 0), (-1, 0), 12),
                        ("SPACEAFTER", (0, 0), (-1, -1), 12),
                    ]
                )
            )
            story.append(image_table)
        else:
            story.append(Paragraph("No hay imágenes disponibles.", normal_style))

        story.append(PageBreak())

        #####################################################################################################
        # Nueva página para el contenido

        if formulario_proyecto.fondos:
            story.append(
                Paragraph("<strong>Postulación a Fondos</strong>", red_bold_style)
            )
            drawing = Drawing()
            line_width = 550
            line = Line(0, 0, line_width, 0)
            line.strokeColor = red
            line.strokeWidth = 1
            drawing.add(line)
            drawing.height = 10
            drawing.width = line_width
            story.append(drawing)

            story.append(Spacer(1, 12))
            fondos = formulario_proyecto.fondos_proyecto.first()
            story.append(
                Paragraph(
                    "<strong>Technology Readiness Levels (TRL):</strong> <font color='%s'>%s</font>"
                    % (dark_gray_color, fondos.trl_id.id),
                    normal_style,
                )
            )
            story.append(Spacer(1, 12))
            story.append(
                Paragraph(
                    "<strong>Problema/Oportunidad:</strong> <font color='%s'>%s</font>"
                    % (dark_gray_color, fondos.problema_oportunidad),
                    normal_style,
                )
            )
            story.append(Spacer(1, 12))
            story.append(
                Paragraph(
                    "<strong>Solución Innovadora:</strong> <font color='%s'>%s</font>"
                    % (dark_gray_color, fondos.solucion_innovadora),
                    normal_style,
                )
            )
            story.append(Spacer(1, 12))
            story.append(
                Paragraph(
                    "<strong>Plan de Trabajo:</strong> <font color='%s'>%s</font>"
                    % (dark_gray_color, fondos.plan_trabajo),
                    normal_style,
                )
            )
            story.append(Spacer(1, 12))
            story.append(
                Paragraph(
                    "<strong>Potencial de Comercialización:</strong> <font color='%s'>%s</font>"
                    % (dark_gray_color, fondos.potencial_comercializacion),
                    normal_style,
                )
            )

        doc.build(story)
        return response

    def get(self, request, *args, **kwargs):
        cod_id = kwargs.get("pk")
        formulario_proyecto = get_object_or_404(FormularioProyectoFabrica, id=cod_id)
        return self.generate_pdf(formulario_proyecto)


###########################
### ## Generate PTT ###
###########################


from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image as PilImage


class GeneratePptFabLabView(LoginRequiredMixin, View):
    def generate_ppt(self, formulario_proyecto):
        # Crear una presentación en memoria
        presentation = Presentation()

        # Agregar una diapositiva con un título y un subtítulo
        slide_layout = presentation.slide_layouts[0]  # Título y subtítulo
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "FABLAB INACAP"
        subtitle.text = "Fichas de Proyecto"

        # Agregar una diapositiva con imágenes del proyecto
        if formulario_proyecto.images.exists():
            slide_layout = presentation.slide_layouts[5]  # Diapositiva en blanco
            slide = presentation.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = "Detalle del Proyecto"

            # Configurar el tamaño y la posición de las imágenes
            left_margin = Inches(0.5)
            top_margin = Inches(1.3)
            image_width = Inches(2.5)
            image_height = Inches(2.5)
            space_between_images = Inches(0.2)  # Espacio entre imágenes
            max_images_per_row = 2

            # Añadir título y TRL al lado derecho de las imágenes
            right_text_box = slide.shapes.add_textbox(
                Inches(6), Inches(1.2), Inches(4), Inches(3)
            )
            text_frame = right_text_box.text_frame
            p = text_frame.add_paragraph()
            p.text = f"Nombre propuesta:\n{formulario_proyecto.nombre_propuesta}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.space_after = Pt(14)

            p = text_frame.add_paragraph()
            p.text = f"TRL:{formulario_proyecto.trl_id}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.space_after = Pt(14)

            # Añadir imágenes
            row = 0
            col = 0
            for image in formulario_proyecto.images.all():
                img_stream = BytesIO(image.image.read())

                # Convertir WEBP a PNG
                with PilImage.open(img_stream) as img:
                    img = img.convert("RGB")  # Convertir a RGB si es necesario
                    png_io = BytesIO()
                    img.save(png_io, format="PNG")
                    png_io.seek(0)

                    # Calcular posición de la imagen
                    left = left_margin + col * (image_width + space_between_images)
                    top = top_margin + row * (image_height + space_between_images)

                    # Agregar imagen al PPT
                    slide.shapes.add_picture(
                        png_io, left, top, width=image_width, height=image_height
                    )

                    # Ajustar fila y columna para la siguiente imagen
                    col += 1
                    if col >= max_images_per_row:
                        col = 0
                        row += 1

            # Añadir información de docentes y alumnos
            text_box = slide.shapes.add_textbox(
                Inches(6), Inches(5), Inches(4), Inches(2)
            )
            text_frame = text_box.text_frame
            p = text_frame.add_paragraph()
            p.text = "Docentes:"
            p.font.size = Pt(12)
            p.font.bold = True
            p.space_after = Pt(4)

            # Lista de docentes
            docentes_list = "\n".join(
                f"• {docente.nombre}" for docente in formulario_proyecto.docentes.all()
            )
            p = text_frame.add_paragraph()
            p.text = docentes_list
            p.font.size = Pt(12)
            p.space_after = Pt(6)

            p = text_frame.add_paragraph()
            p.text = f"Alumnos Participantes: {formulario_proyecto.alumnos}"
            p.font.size = Pt(12)
            p.font.bold = True

        # Agregar diapositiva con el problema
        slide_layout = presentation.slide_layouts[1]  # Título y contenido
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Problema"
        content.text = formulario_proyecto.problema
        content.text_frame.paragraphs[0].font.size = Pt(14)

        # Agregar diapositiva con la solución
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Solución"
        content.text = formulario_proyecto.solucion
        content.text_frame.paragraphs[0].font.size = Pt(14)

        # Guardar la presentación en un objeto BytesIO
        byte_io = BytesIO()
        presentation.save(byte_io)
        byte_io.seek(0)

        # Crear la respuesta HTTP
        response = HttpResponse(
            byte_io,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        response["Content-Disposition"] = 'attachment; filename="presentacion.pptx"'

        return response

    def get(self, request, *args, **kwargs):
        cod_id = kwargs.get("pk")
        formulario_proyecto = get_object_or_404(FormularioProyectoFabLab, id=cod_id)
        return self.generate_ppt(formulario_proyecto)
